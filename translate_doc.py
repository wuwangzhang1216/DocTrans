import os
import json
from typing import Dict, List, Optional, Tuple
from pathlib import Path
from google import genai
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
import re
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# PyMuPDF for improved PDF handling
try:
    import fitz  # PyMuPDF
except ImportError:
    import pymupdf as fitz

# Document processing libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
import platform

@dataclass
class TextBlock:
    """Data class to store text block information"""
    bbox: tuple  # (x0, y0, x1, y1)
    text: str
    font: str = None
    font_size: float = None
    flags: int = 0
    color: int = 0
    block_no: int = 0
    block_type: int = 0

class DocumentTranslator:
    """
    An improved document translator that better preserves PDF layout
    using PyMuPDF's advanced features for overlay and redaction.
    Supports PDF, PPTX, DOCX, and TXT files.
    """
    
    def __init__(self, api_key: str = None, model: str = "gemini-2.5-flash-lite-preview-09-2025", max_workers: int = 256):
        """
        Initialize the translator with Google Generative AI.

        Args:
            api_key: Google API key (optional, will use GEMINI_API_KEY env var if not provided)
            model: Gemini model to use for translation (default: gemini-2.5-flash-lite-preview-09-2025)
            max_workers: Maximum total number of parallel workers (default: 256)
        """
        if api_key:
            os.environ['GEMINI_API_KEY'] = api_key
        self.client = genai.Client()
        self.model = model
        self.max_workers = max_workers
        self.max_concurrent_pages = 16
        self.max_workers_per_page = 64

        # Compile regex patterns for math/LaTeX detection
        self._compile_math_patterns()

    def _compile_math_patterns(self):
        """Compile regex patterns for detecting LaTeX and mathematical expressions."""
        # LaTeX display math: $$...$$ or \[...\]
        self.display_math_patterns = [
            re.compile(r'\$\$.*?\$\$', re.DOTALL),  # $$...$$
            re.compile(r'\\\[.*?\\\]', re.DOTALL),   # \[...\]
        ]

        # LaTeX inline math: $...$ or \(...\)
        self.inline_math_patterns = [
            re.compile(r'(?<!\$)\$(?!\$)([^\$]+?)\$(?!\$)'),  # $...$ (not $$)
            re.compile(r'\\\(.*?\\\)'),  # \(...\)
        ]

        # Common LaTeX commands and environments
        self.latex_command_patterns = [
            re.compile(r'\\[a-zA-Z]+\{[^}]*\}'),  # \command{...}
            re.compile(r'\\begin\{[^}]+\}.*?\\end\{[^}]+\}', re.DOTALL),  # \begin{env}...\end{env}
            re.compile(r'\\(?:frac|sqrt|sum|int|lim|prod|infty|alpha|beta|gamma|delta|theta|lambda|sigma|omega|Delta|Gamma|Theta|Lambda|Sigma|Omega)\b'),  # Common math commands
        ]

        # Mathematical symbols and operators
        self.math_symbol_patterns = [
            re.compile(r'[∫∑∏∂∇∆√∞≈≠≤≥±×÷∈∉⊂⊃⊆⊇∪∩∧∨¬∀∃→←↔⇒⇐⇔]'),  # Unicode math symbols
            re.compile(r'[αβγδεζηθικλμνξοπρστυφχψω]'),  # Greek letters
            re.compile(r'[ΑΒΓΔΕΖΗΘΙΚΛΜΝΞΟΠΡΣΤΥΦΧΨΩ]'),  # Capital Greek letters
        ]

    def _protect_math_expressions(self, text: str) -> Tuple[str, Dict[str, str]]:
        """
        Replace math expressions with placeholders to prevent translation.

        Args:
            text: Text containing potential math expressions

        Returns:
            Tuple of (protected_text, placeholder_map)
            - protected_text: Text with math expressions replaced by placeholders
            - placeholder_map: Dictionary mapping placeholders to original expressions
        """
        if not text or not text.strip():
            return text, {}

        protected_text = text
        placeholder_map = {}
        placeholder_counter = 0

        # Protect display math first (highest priority)
        for pattern in self.display_math_patterns:
            matches = pattern.finditer(protected_text)
            for match in matches:
                math_expr = match.group(0)
                placeholder = f"__MATH_DISPLAY_{placeholder_counter}__"
                placeholder_map[placeholder] = math_expr
                protected_text = protected_text.replace(math_expr, placeholder, 1)
                placeholder_counter += 1

        # Protect inline math
        for pattern in self.inline_math_patterns:
            matches = pattern.finditer(protected_text)
            for match in matches:
                math_expr = match.group(0)
                placeholder = f"__MATH_INLINE_{placeholder_counter}__"
                placeholder_map[placeholder] = math_expr
                protected_text = protected_text.replace(math_expr, placeholder, 1)
                placeholder_counter += 1

        # Protect LaTeX commands and environments
        for pattern in self.latex_command_patterns:
            matches = pattern.finditer(protected_text)
            for match in matches:
                latex_expr = match.group(0)
                # Skip if already protected
                if latex_expr.startswith('__MATH_'):
                    continue
                placeholder = f"__LATEX_CMD_{placeholder_counter}__"
                placeholder_map[placeholder] = latex_expr
                protected_text = protected_text.replace(latex_expr, placeholder, 1)
                placeholder_counter += 1

        # Protect mathematical symbols (only if they appear in mathematical context)
        # We're conservative here - only protect if surrounded by numbers or other math
        for pattern in self.math_symbol_patterns:
            matches = pattern.finditer(protected_text)
            for match in matches:
                symbol = match.group(0)
                start = match.start()
                end = match.end()

                # Check context: is it surrounded by numbers, operators, or other math?
                before = protected_text[max(0, start-1):start] if start > 0 else ''
                after = protected_text[end:min(len(protected_text), end+1)] if end < len(protected_text) else ''

                # Protect if in mathematical context
                if (before and (before.isdigit() or before in '+-*/=()[]{}.,<> ')) or \
                   (after and (after.isdigit() or after in '+-*/=()[]{}.,<> ')):
                    placeholder = f"__MATH_SYM_{placeholder_counter}__"
                    placeholder_map[placeholder] = symbol
                    protected_text = protected_text[:start] + placeholder + protected_text[end:]
                    placeholder_counter += 1

        return protected_text, placeholder_map

    def _restore_math_expressions(self, text: str, placeholder_map: Dict[str, str]) -> str:
        """
        Restore original math expressions from placeholders.

        Args:
            text: Translated text with placeholders
            placeholder_map: Dictionary mapping placeholders to original expressions

        Returns:
            Text with math expressions restored
        """
        if not placeholder_map:
            return text

        restored_text = text

        # Restore in reverse order to handle nested cases correctly
        for placeholder, original in placeholder_map.items():
            restored_text = restored_text.replace(placeholder, original)

        return restored_text

    def _calculate_workers_per_page(self, total_pages: int) -> tuple:
        """
        Calculate dynamic worker allocation based on total pages.

        Args:
            total_pages: Total number of pages to process

        Returns:
            Tuple of (workers_per_page, concurrent_pages)
        """
        # Limit concurrent pages to max_concurrent_pages
        concurrent_pages = min(total_pages, self.max_concurrent_pages)

        # Calculate workers per page: total_workers / concurrent_pages
        # But cap at max_workers_per_page
        workers_per_page = min(
            self.max_workers_per_page,
            self.max_workers // concurrent_pages
        )

        return workers_per_page, concurrent_pages

    def translate_text(self, text: str, target_language: str,
                      source_language: str = "auto") -> str:
        """
        Translate text using Google Generative AI.

        Args:
            text: Text to translate
            target_language: Target language for translation
            source_language: Source language (default: auto-detect)

        Returns:
            Translated text
        """
        if not text or not text.strip():
            return text

        try:
            # Protect math expressions before translation
            protected_text, math_map = self._protect_math_expressions(text)

            prompt = (
                f"You are a professional technical translator. Translate into {target_language} with precise domain terminology and preserve formatting exactly.\n\n"
                f"Translate to {target_language} with native, accurate, technical wording.\n"
                "Strictly preserve original formatting and layout: line breaks, indentation, spacing, bullet/numbered lists (markers and levels), tables, and code blocks.\n"
                "Do not add explanations. Do not change capitalization of proper nouns.\n"
                "Do not translate code, CLI commands, file paths, API names, or placeholders like {var}, <tag>, {{braces}}, [1], %s, or ${VAR}.\n"
                "CRITICAL: Do NOT translate LaTeX expressions, mathematical formulas, equations, or any placeholders starting with __MATH_ or __LATEX_. Keep them EXACTLY as they appear.\n"
                "Keep URLs and IDs unchanged.\n\n"
                "Text to translate:\n"
                f"{protected_text}"
            )

            response = self.client.models.generate_content(
                model=self.model,
                contents=prompt
            )

            # Restore math expressions after translation
            translated_text = response.text.strip()
            final_text = self._restore_math_expressions(translated_text, math_map)

            return final_text

        except Exception as e:
            print(f"Translation error: {str(e)}")
            return text

    # ============= PPTX Translation Methods =============
    
    def translate_slide_content(self, slide_data: tuple) -> dict:
        """
        Translate content of a single slide. Helper method for parallel processing.

        Args:
            slide_data: Tuple containing (slide_index, slide_content, target_language, api_key, model)

        Returns:
            Dictionary with translated content for the slide
        """
        slide_idx, shapes_data, target_language, api_key, model = slide_data

        try:
            # Create a new client instance for thread safety
            if api_key:
                os.environ['GEMINI_API_KEY'] = api_key
            client = genai.Client()
            translated_shapes = []

            def translate_runs_with_context(runs: List[str]) -> List[str]:
                if not runs:
                    return runs
                user_prompt = (
                    f"You are a professional technical translator. Translate into {target_language} precisely and preserve formatting.\n\n"
                    "Translate the following paragraph into "
                    f"{target_language} with native, accurate technical language.\n"
                    "Rules:\n"
                    "- Consider all runs together for context, but return the translation per run.\n"
                    "- Keep the number of runs identical and in the same order.\n"
                    "- Do not merge or split runs; do not drop whitespace.\n"
                    "- Preserve placeholders, code, variables, URLs, IDs unchanged.\n"
                    "- Output ONLY a JSON array of strings of the same length; no backticks, no extra text.\n\n"
                    f"Runs: {json.dumps(runs, ensure_ascii=False)}"
                )

                resp = client.models.generate_content(
                    model=model,
                    contents=user_prompt
                )
                content = resp.text.strip()
                try:
                    data = json.loads(content)
                    if isinstance(data, list) and len(data) == len(runs) and all(isinstance(x, str) for x in data):
                        return data
                except Exception:
                    pass
                return runs

            for shape_idx, shape_info in enumerate(shapes_data):
                shape_type, content = shape_info

                if shape_type == 'text':
                    if content and content.strip():
                        prompt = (
                            f"You are a professional technical translator. Translate into {target_language} and preserve formatting.\n\n"
                            f"Translate to {target_language} with native, accurate, technical wording.\n"
                            "Strictly preserve original line breaks, indentation, and list markers.\n"
                            "Only return the translated text.\n\n"
                            f"Text to translate:\n{content}"
                        )

                        response = client.models.generate_content(
                            model=model,
                            contents=prompt
                        )
                        translated_shapes.append((shape_idx, response.text.strip()))
                    else:
                        translated_shapes.append((shape_idx, content))

                elif shape_type == 'table':
                    # content is nested: rows -> cells -> paragraphs -> runs
                    translated_table = []
                    for row_idx, row_cells in enumerate(content):
                        translated_row_cells = []
                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                            translated_cell_paragraphs = []
                            for para_idx, runs in enumerate(cell_paragraphs):
                                if runs and any(r.strip() for r in runs):
                                    translated_runs = translate_runs_with_context(runs)
                                else:
                                    translated_runs = runs
                                translated_cell_paragraphs.append(translated_runs)
                            translated_row_cells.append(translated_cell_paragraphs)
                        translated_table.append(translated_row_cells)
                    translated_shapes.append((shape_idx, translated_table))

                elif shape_type == 'paragraphs':
                    translated_paragraphs = []
                    for para_idx, runs_data in enumerate(content):
                        translated_runs = translate_runs_with_context(runs_data)
                        translated_paragraphs.append((
                            para_idx,
                            [(i, t) for i, t in enumerate(translated_runs)]
                        ))
                    translated_shapes.append((shape_idx, translated_paragraphs))

            return {'slide_idx': slide_idx, 'translated_shapes': translated_shapes}

        except Exception as e:
            print(f"Translation error on slide {slide_idx + 1}: {str(e)}")
            return {'slide_idx': slide_idx, 'error': str(e)}

    def translate_pptx(self, input_path: str, output_path: str,
                      target_language: str) -> bool:
        """
        Translate PowerPoint presentation.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to save translated PPTX file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            prs = Presentation(input_path)

            # Extract content from all slides for parallel processing
            slides_data = []
            for slide_idx, slide in enumerate(prs.slides):
                shapes_data = []

                for shape_idx, shape in enumerate(slide.shapes):
                    # Handle tables first, then text frames to preserve bullets/styles
                    if shape.has_table:
                        table_data = []  # rows -> cells -> paragraphs -> runs
                        table = shape.table
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_paragraphs = []
                                try:
                                    paragraphs = list(cell.text_frame.paragraphs)
                                except Exception:
                                    paragraphs = []
                                for paragraph in paragraphs:
                                    runs_data = [run.text for run in paragraph.runs]
                                    cell_paragraphs.append(runs_data)
                                row_cells.append(cell_paragraphs)
                            table_data.append(row_cells)
                        shapes_data.append(('table', table_data))
                    elif shape.has_text_frame:
                        paragraphs_data = []
                        for paragraph in shape.text_frame.paragraphs:
                            runs_data = [run.text for run in paragraph.runs]
                            paragraphs_data.append(runs_data)
                        shapes_data.append(('paragraphs', paragraphs_data))
                    elif hasattr(shape, "text") and shape.text:
                        # Fallback for shapes exposing only plain text
                        shapes_data.append(('text', shape.text))
                    else:
                        shapes_data.append(('empty', None))

                if shapes_data:  # Only process slides that have content
                    api_key = os.environ.get('GEMINI_API_KEY', None)
                    slides_data.append((slide_idx, shapes_data, target_language, api_key, self.model))

            if not slides_data:
                print("No text content found in PPTX")
                prs.save(output_path)
                return True

            total_slides = len(slides_data)
            workers_per_slide, concurrent_slides = self._calculate_workers_per_page(total_slides)

            print(f"Processing {total_slides} slides with dynamic allocation:")
            print(f"  - Processing {concurrent_slides} slides concurrently")
            print(f"  - Using up to {workers_per_slide} workers per slide")
            print(f"  - Total workers: {workers_per_slide * concurrent_slides}")

            # Process slides in parallel with dynamic allocation
            translated_results = {}
            with ThreadPoolExecutor(max_workers=workers_per_slide * concurrent_slides) as executor:
                future_to_slide = {
                    executor.submit(self.translate_slide_content, slide_data): slide_data[0]
                    for slide_data in slides_data
                }

                for future in concurrent.futures.as_completed(future_to_slide):
                    slide_idx = future_to_slide[future]
                    try:
                        result = future.result()
                        if 'error' not in result:
                            translated_results[slide_idx] = result['translated_shapes']
                    except Exception as e:
                        print(f"Slide {slide_idx + 1} translation failed: {str(e)}")

            # Apply translations back to the presentation
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx in translated_results:
                    shape_translations = translated_results[slide_idx]
                    shape_translations.sort(key=lambda x: x[0])  # Sort by shape index

                    for shape_trans_idx, (shape_idx, translated_content) in enumerate(shape_translations):
                        if shape_idx < len(list(slide.shapes)):
                            shape = list(slide.shapes)[shape_idx]

                            # Prefer text_frame updates (preserves bullets and run formatting)
                            if shape.has_text_frame and isinstance(translated_content, list):
                                for para_trans in translated_content:
                                    if isinstance(para_trans, tuple) and len(para_trans) == 2:
                                        para_idx, translated_runs = para_trans
                                        if para_idx < len(list(shape.text_frame.paragraphs)):
                                            paragraph = list(shape.text_frame.paragraphs)[para_idx]
                                            for run_trans in translated_runs:
                                                if isinstance(run_trans, tuple) and len(run_trans) == 2:
                                                    run_idx, run_text = run_trans
                                                    if run_idx < len(list(paragraph.runs)):
                                                        list(paragraph.runs)[run_idx].text = run_text
                            elif shape.has_table and isinstance(translated_content, list):
                                # Apply nested table translations: rows -> cells -> paragraphs -> runs
                                table = shape.table
                                for row_idx, row_cells in enumerate(translated_content):
                                    if row_idx < len(list(table.rows)):
                                        row = list(table.rows)[row_idx]
                                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                                            if cell_idx < len(list(row.cells)):
                                                cell = list(row.cells)[cell_idx]
                                                for para_idx, translated_runs in enumerate(cell_paragraphs):
                                                    if para_idx < len(list(cell.text_frame.paragraphs)):
                                                        paragraph = list(cell.text_frame.paragraphs)[para_idx]
                                                        for run_idx, run_text in enumerate(translated_runs):
                                                            if run_idx < len(list(paragraph.runs)):
                                                                list(paragraph.runs)[run_idx].text = run_text
                            elif hasattr(shape, "text") and isinstance(translated_content, str):
                                shape.text = translated_content

            # Save translated presentation
            prs.save(output_path)
            print(f"[SUCCESS] Translated PPTX saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating PPTX: {str(e)}")
            return False

    # ============= DOCX Translation Methods =============
    
    def translate_paragraph_batch(self, batch_data: tuple) -> dict:
        """
        Translate a batch of paragraphs (run-aligned). Helper for parallel processing.

        Args:
            batch_data: Tuple containing (batch_index, paragraphs_runs_data, target_language, api_key, model)

        Returns:
            Dictionary with translated paragraphs for the batch
        """
        batch_idx, paragraphs_runs_data, target_language, api_key, model = batch_data

        try:
            if api_key:
                os.environ['GEMINI_API_KEY'] = api_key
            client = genai.Client()
            results = []

            def translate_runs(runs: List[str]) -> List[str]:
                if not runs:
                    return runs
                prompt = (
                    f"You are a professional technical translator. Translate into {target_language} precisely and preserve formatting.\n\n"
                    f"Translate into {target_language} with native, accurate, technical wording.\n"
                    "Consider all runs together for context, but return the translation per run. Keep the same number of runs and order.\n"
                    "Do not merge or split runs; preserve whitespace.\n"
                    "Preserve placeholders, code, variables, URLs, and IDs unchanged.\n"
                    "Respond with ONLY a JSON array of strings of the same length.\n\n"
                    f"Runs: {json.dumps(runs, ensure_ascii=False)}"
                )
                resp = client.models.generate_content(
                    model=model,
                    contents=prompt
                )
                content = resp.text.strip()
                try:
                    data = json.loads(content)
                    if isinstance(data, list) and len(data) == len(runs) and all(isinstance(x, str) for x in data):
                        return data
                except Exception:
                    pass
                return runs

            for item in paragraphs_runs_data:
                para_idx, runs = item
                if runs and any(t.strip() for t in runs):
                    translated_runs = translate_runs(runs)
                else:
                    translated_runs = runs
                results.append((para_idx, [(i, t) for i, t in enumerate(translated_runs)]))

            return {"batch_idx": batch_idx, "translated_paragraphs": results}

        except Exception as e:
            print(f"Translation error in batch {batch_idx}: {str(e)}")
            return {"batch_idx": batch_idx, "error": str(e)}

    def translate_table_batch(self, batch_data: tuple) -> dict:
        """
        Translate a batch of table paragraph runs for DOCX.

        Args:
            batch_data: Tuple containing (batch_index, table_para_runs_data, target_language, api_key, model)

        Returns:
            Dictionary with translated table paragraphs for the batch
        """
        batch_idx, table_para_runs_data, target_language, api_key, model = batch_data

        try:
            if api_key:
                os.environ['GEMINI_API_KEY'] = api_key
            client = genai.Client()
            results = []

            def translate_runs(runs: List[str]) -> List[str]:
                if not runs:
                    return runs
                prompt = (
                    f"You are a professional technical translator. Translate into {target_language} precisely and preserve formatting.\n\n"
                    f"Translate into {target_language} with native, accurate, technical wording.\n"
                    "Consider all runs together for context, but return the translation per run. Keep the same number of runs and order.\n"
                    "Do not merge or split runs; preserve whitespace.\n"
                    "Preserve placeholders, code, variables, URLs, and IDs unchanged.\n"
                    "Respond with ONLY a JSON array of strings of the same length.\n\n"
                    f"Runs: {json.dumps(runs, ensure_ascii=False)}"
                )
                resp = client.models.generate_content(
                    model=model,
                    contents=prompt
                )
                content = resp.text.strip()
                try:
                    data = json.loads(content)
                    if isinstance(data, list) and len(data) == len(runs) and all(isinstance(x, str) for x in data):
                        return data
                except Exception:
                    pass
                return runs

            for item in table_para_runs_data:
                key, runs = item  # key identifies (table_idx, row_idx, cell_idx, para_idx)
                if runs and any(t.strip() for t in runs):
                    translated_runs = translate_runs(runs)
                else:
                    translated_runs = runs
                results.append((key, [(i, t) for i, t in enumerate(translated_runs)]))

            return {"batch_idx": batch_idx, "translated_table_paragraphs": results}

        except Exception as e:
            print(f"Translation error in table batch {batch_idx}: {str(e)}")
            return {"batch_idx": batch_idx, "error": str(e)}

    def translate_docx(self, input_path: str, output_path: str,
                      target_language: str) -> bool:
        """
        Translate Word document.

        Args:
            input_path: Path to input DOCX file
            output_path: Path to save translated DOCX file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            doc = Document(input_path)

            # Gather paragraph runs for translation (document body)
            body_paragraphs = []  # list of (para_idx, [run_texts])
            for para_idx, paragraph in enumerate(doc.paragraphs):
                runs_texts = [r.text for r in paragraph.runs]
                body_paragraphs.append((para_idx, runs_texts))

            # Gather table cell paragraph runs for translation
            # Index by (table_idx, row_idx, cell_idx, para_idx)
            table_para_runs = []
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, paragraph in enumerate(cell.paragraphs):
                            runs_texts = [r.text for r in paragraph.runs]
                            table_para_runs.append(((table_idx, row_idx, cell_idx, para_idx), runs_texts))

            # Calculate total items for dynamic worker allocation
            total_items = len(body_paragraphs) + len(table_para_runs)
            workers_allocated, _ = self._calculate_workers_per_page(max(1, total_items // 10))  # Treat every ~10 items as a "page"

            print(f"Processing DOCX with dynamic allocation:")
            print(f"  - Body paragraphs: {len(body_paragraphs)}")
            print(f"  - Table paragraphs: {len(table_para_runs)}")
            print(f"  - Using up to {workers_allocated} workers")

            # Process body paragraphs in batches
            paragraph_batches = []
            if body_paragraphs:
                batch_size = max(1, len(body_paragraphs) // workers_allocated) if workers_allocated > 1 else len(body_paragraphs)
                for i in range(0, len(body_paragraphs), batch_size):
                    batch = body_paragraphs[i:i + batch_size]
                    api_key = os.environ.get('GEMINI_API_KEY', None)
                    paragraph_batches.append((len(paragraph_batches), batch, target_language, api_key, self.model))

            translated_body = {}
            if paragraph_batches:
                print(f"Processing {len(paragraph_batches)} paragraph batches with {min(workers_allocated, len(paragraph_batches))} parallel workers...")
                with ThreadPoolExecutor(max_workers=min(workers_allocated, len(paragraph_batches))) as executor:
                    future_to_batch = {
                        executor.submit(self.translate_paragraph_batch, batch): batch[0]
                        for batch in paragraph_batches
                    }

                    for future in concurrent.futures.as_completed(future_to_batch):
                        batch_idx = future_to_batch[future]
                        try:
                            result = future.result()
                            if 'error' not in result:
                                for para_idx, translated_runs in result['translated_paragraphs']:
                                    translated_body[para_idx] = translated_runs
                        except Exception as e:
                            print(f"Paragraph batch {batch_idx} translation failed: {str(e)}")

            # Process table paragraphs in batches
            table_batches = []
            if table_para_runs:
                batch_size = max(1, len(table_para_runs) // workers_allocated) if workers_allocated > 1 else len(table_para_runs)
                for i in range(0, len(table_para_runs), batch_size):
                    batch = table_para_runs[i:i + batch_size]
                    api_key = os.environ.get('GEMINI_API_KEY', None)
                    table_batches.append((len(table_batches), batch, target_language, api_key, self.model))

            translated_table = {}
            if table_batches:
                print(f"Processing {len(table_batches)} table batches with {min(workers_allocated, len(table_batches))} parallel workers...")
                with ThreadPoolExecutor(max_workers=min(workers_allocated, len(table_batches))) as executor:
                    future_to_batch = {
                        executor.submit(self.translate_table_batch, batch): batch[0]
                        for batch in table_batches
                    }

                    for future in concurrent.futures.as_completed(future_to_batch):
                        batch_idx = future_to_batch[future]
                        try:
                            result = future.result()
                            if 'error' not in result:
                                for key, translated_runs in result['translated_table_paragraphs']:
                                    translated_table[key] = translated_runs
                        except Exception as e:
                            print(f"Table batch {batch_idx} translation failed: {str(e)}")

            # Apply body paragraph translations (preserve run styles like underline/bold)
            for para_idx, paragraph in enumerate(doc.paragraphs):
                if para_idx in translated_body:
                    for run_idx, run_text in translated_body[para_idx]:
                        if run_idx < len(paragraph.runs):
                            paragraph.runs[run_idx].text = run_text

            # Apply table translations (preserve run styles)
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, paragraph in enumerate(cell.paragraphs):
                            key = (table_idx, row_idx, cell_idx, para_idx)
                            if key in translated_table:
                                for run_idx, run_text in translated_table[key]:
                                    if run_idx < len(paragraph.runs):
                                        paragraph.runs[run_idx].text = run_text

            # Save translated document
            doc.save(output_path)
            print(f"[SUCCESS] Translated DOCX saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating DOCX: {str(e)}")
            return False

    # ============= PDF Translation Methods (Improved) =============

    def translate_text_block(self, block_data: tuple) -> dict:
        """
        Translate a single text block. Helper method for parallel processing.

        Args:
            block_data: Tuple containing (block_index, block_text, target_language, api_key, model)

        Returns:
            Dictionary with translated content for the block
        """
        block_idx, block_text, target_language, api_key, model = block_data

        try:
            # Create a new client instance for thread safety
            if api_key:
                os.environ['GEMINI_API_KEY'] = api_key
            client = genai.Client()

            # Protect math expressions before translation
            protected_text, math_map = self._protect_math_expressions(block_text)

            prompt = (
                f"You are a professional technical translator. Translate into {target_language} with precise domain terminology and preserve formatting exactly.\n\n"
                f"Translate to {target_language} with native, accurate, technical wording.\n"
                "Strictly preserve original formatting and layout: line breaks, indentation, spacing, bullet/numbered lists (markers and levels), tables, and code blocks.\n"
                "Do not add explanations. Do not change capitalization of proper nouns.\n"
                "Do not translate code, CLI commands, file paths, API names, or placeholders like {var}, <tag>, {{braces}}, [1], %s, or ${VAR}.\n"
                "CRITICAL: Do NOT translate LaTeX expressions, mathematical formulas, equations, or any placeholders starting with __MATH_ or __LATEX_. Keep them EXACTLY as they appear.\n"
                "Keep URLs and IDs unchanged.\n\n"
                "Text to translate:\n"
                f"{protected_text}"
            )

            response = client.models.generate_content(
                model=model,
                contents=prompt
            )

            # Restore math expressions after translation
            translated_text = response.text.strip()
            final_text = self._restore_math_expressions(translated_text, math_map)

            return {'block_idx': block_idx, 'translated_text': final_text}

        except Exception as e:
            print(f"Translation error on block {block_idx}: {str(e)}")
            return {'block_idx': block_idx, 'translated_text': block_text, 'error': str(e)}

    def translate_pdf_with_overlay(self, input_path: str, output_path: str,
                                  target_language: str) -> bool:
        """
        Translate PDF using overlay method to preserve layout.
        This method adds translated text over the original text without duplicating pages.

        Args:
            input_path: Path to input PDF file
            output_path: Path to save translated PDF file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            # Open the PDF document directly (modify in place)
            doc = fitz.open(input_path)

            total_pages = len(doc)
            workers_per_page, concurrent_pages = self._calculate_workers_per_page(total_pages)

            print(f"Processing {total_pages} pages with dynamic allocation:")
            print(f"  - Processing {concurrent_pages} pages concurrently")
            print(f"  - Using up to {workers_per_page} workers per page")
            print(f"  - Total workers: {workers_per_page * concurrent_pages}")

            # Extract all pages data first
            pages_data = []
            for page_num in range(total_pages):
                page = doc[page_num]
                blocks = page.get_text("dict")
                page_blocks = []

                for block in blocks.get("blocks", []):
                    if block.get("type") == 0:  # Text block
                        bbox = fitz.Rect(block["bbox"])

                        # Extract text from the block
                        block_text = ""
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                if span.get("text"):
                                    block_text += span["text"] + " "

                        if block_text.strip():
                            # Get font information from the first span
                            font_info = None
                            for line in block.get("lines", []):
                                for span in line.get("spans", []):
                                    font_info = span
                                    break
                                if font_info:
                                    break

                            page_blocks.append({
                                'bbox': bbox,
                                'text': block_text.strip(),
                                'font_size': font_info.get("size", 11) if font_info else 11,
                                'font_color': font_info.get("color", 0) if font_info else 0
                            })

                pages_data.append({
                    'page_num': page_num,
                    'blocks': page_blocks
                })

            # Process pages with dynamic worker allocation
            api_key = os.environ.get('GEMINI_API_KEY', None)

            for i in range(0, total_pages, concurrent_pages):
                batch_pages = pages_data[i:i + concurrent_pages]
                print(f"\nProcessing pages {i + 1}-{min(i + concurrent_pages, total_pages)}...")

                # Prepare all blocks from this batch of pages for parallel translation
                blocks_to_translate = []
                for page_data in batch_pages:
                    for block_idx, block in enumerate(page_data['blocks']):
                        global_idx = (page_data['page_num'], block_idx)
                        blocks_to_translate.append((
                            global_idx,
                            block['text'],
                            target_language,
                            api_key,
                            self.model
                        ))

                # Translate all blocks in parallel with workers_per_page
                translated_blocks = {}
                if blocks_to_translate:
                    with ThreadPoolExecutor(max_workers=workers_per_page) as executor:
                        future_to_block = {
                            executor.submit(self.translate_text_block, block_data): block_data[0]
                            for block_data in blocks_to_translate
                        }

                        for future in concurrent.futures.as_completed(future_to_block):
                            block_idx = future_to_block[future]
                            try:
                                result = future.result()
                                if 'error' not in result:
                                    translated_blocks[block_idx] = result['translated_text']
                            except Exception as e:
                                print(f"Block {block_idx} translation failed: {str(e)}")

                # Apply translations to pages (in place overlay)
                for page_data in batch_pages:
                    page_num = page_data['page_num']
                    page = doc[page_num]

                    # Apply translated blocks
                    for block_idx, block in enumerate(page_data['blocks']):
                        global_idx = (page_num, block_idx)

                        # Draw white rectangle to cover original text
                        page.draw_rect(block['bbox'], color=(1, 1, 1), fill=(1, 1, 1))

                        # Insert translated text if available
                        if global_idx in translated_blocks:
                            translated_text = translated_blocks[global_idx]
                            color_rgb = self._int_to_rgb(block['font_color'])

                            # Use htmlbox for reliable rendering (handles all languages and auto-wraps)
                            html = f'<span style="font-size:{block["font_size"]}pt; color:rgb{color_rgb};">{translated_text}</span>'
                            page.insert_htmlbox(
                                block['bbox'],
                                html,
                                css="body { margin: 0; padding: 2px; }"
                            )

            # Save the translated PDF with maximum compression
            doc.save(
                output_path,
                garbage=4,        # Most aggressive garbage collection
                deflate=True,     # Compress streams
                clean=True,       # Remove duplicate objects
                pretty=False,     # No pretty printing
                no_new_id=True,   # Don't change file ID
                linear=False      # No web optimization (smaller file)
            )
            doc.close()

            print(f"\n[SUCCESS] Translated PDF saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating PDF with overlay: {str(e)}")
            return False

    def translate_pdf_with_redaction(self, input_path: str, output_path: str,
                                    target_language: str) -> bool:
        """
        Translate PDF using redaction method - better for text-heavy documents.
        This removes original text and replaces it with translated text.

        Args:
            input_path: Path to input PDF file
            output_path: Path to save translated PDF file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            doc = fitz.open(input_path)

            total_pages = len(doc)
            workers_per_page, concurrent_pages = self._calculate_workers_per_page(total_pages)

            print(f"Processing {total_pages} pages with redaction method and dynamic allocation:")
            print(f"  - Processing {concurrent_pages} pages concurrently")
            print(f"  - Using up to {workers_per_page} workers per page")
            print(f"  - Total workers: {workers_per_page * concurrent_pages}")

            # Extract all pages data first
            pages_data = []
            for page_num, page in enumerate(doc):
                blocks = page.get_text("dict", flags=11)
                page_blocks = []

                for block in blocks.get("blocks", []):
                    if block.get("type") == 0:  # Text block
                        block_bbox = fitz.Rect(block["bbox"])

                        # Collect text from all spans in the block
                        block_text = ""
                        first_span = None

                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                if not first_span:
                                    first_span = span
                                text = span.get("text", "")
                                if text:
                                    block_text += text + " "

                        if block_text.strip() and first_span:
                            page_blocks.append({
                                'bbox': block_bbox,
                                'text': block_text.strip(),
                                'font_size': first_span.get("size", 11),
                                'color': first_span.get("color", 0),
                                'flags': first_span.get("flags", 0)
                            })

                pages_data.append({
                    'page_num': page_num,
                    'blocks': page_blocks
                })

            # Process pages with dynamic worker allocation
            api_key = os.environ.get('GEMINI_API_KEY', None)

            for i in range(0, total_pages, concurrent_pages):
                batch_pages = pages_data[i:i + concurrent_pages]
                print(f"\nProcessing pages {i + 1}-{min(i + concurrent_pages, total_pages)}...")

                # Prepare all blocks from this batch of pages for parallel translation
                blocks_to_translate = []
                for page_data in batch_pages:
                    for block_idx, block in enumerate(page_data['blocks']):
                        global_idx = (page_data['page_num'], block_idx)
                        blocks_to_translate.append((
                            global_idx,
                            block['text'],
                            target_language,
                            api_key,
                            self.model
                        ))

                # Translate all blocks in parallel with workers_per_page
                translated_blocks = {}
                if blocks_to_translate:
                    with ThreadPoolExecutor(max_workers=workers_per_page) as executor:
                        future_to_block = {
                            executor.submit(self.translate_text_block, block_data): block_data[0]
                            for block_data in blocks_to_translate
                        }

                        for future in concurrent.futures.as_completed(future_to_block):
                            block_idx = future_to_block[future]
                            try:
                                result = future.result()
                                if 'error' not in result:
                                    translated_blocks[block_idx] = result['translated_text']
                            except Exception as e:
                                print(f"Block {block_idx} translation failed: {str(e)}")

                # Apply translations to pages
                for page_data in batch_pages:
                    page_num = page_data['page_num']
                    page = doc[page_num]

                    # Add redaction annotations for all blocks
                    for block in page_data['blocks']:
                        page.add_redact_annot(block['bbox'])

                    # Apply redactions (removes original text)
                    page.apply_redactions()

                    # Insert translated text
                    for block_idx, block in enumerate(page_data['blocks']):
                        global_idx = (page_num, block_idx)

                        if global_idx in translated_blocks:
                            translated_text = translated_blocks[global_idx]
                            color_rgb = self._int_to_rgb(block['color'])

                            # Use htmlbox for reliable rendering (handles all languages and auto-wraps)
                            html = (
                                f'<span style="font-size:{block["font_size"]}pt; '
                                f'color:rgb{color_rgb}; '
                                f'font-family: sans-serif;">{translated_text}</span>'
                            )
                            page.insert_htmlbox(
                                block['bbox'],
                                html,
                                css="body { margin: 0; padding: 2px; line-height: 1.2; }"
                            )

            # Optimize and save with maximum compression
            doc.save(
                output_path,
                garbage=4,        # Most aggressive garbage collection
                deflate=True,     # Compress streams
                clean=True,       # Remove duplicate objects
                pretty=False,     # No pretty printing
                no_new_id=True,   # Don't change file ID
                linear=False      # No web optimization (smaller file)
            )
            doc.close()

            print(f"\n[SUCCESS] Translated PDF (redaction method) saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error in redaction translation: {str(e)}")
            return False

    def translate_pdf_hybrid(self, input_path: str, output_path: str,
                           target_language: str) -> bool:
        """
        Hybrid approach: Analyzes the PDF and chooses the best method.
        Uses overlay for PDFs with complex backgrounds, redaction for text-heavy docs.
        
        Args:
            input_path: Path to input PDF file
            output_path: Path to save translated PDF file  
            target_language: Target language for translation
            
        Returns:
            Success status
        """
        try:
            # Analyze the PDF to determine best approach
            doc = fitz.open(input_path)
            
            # Check if PDF has images or complex backgrounds
            has_complex_background = False
            total_images = 0
            total_text_blocks = 0
            
            for page in doc:
                # Count images
                image_list = page.get_images()
                total_images += len(image_list)
                
                # Count text blocks
                blocks = page.get_text("dict")
                for block in blocks.get("blocks", []):
                    if block.get("type") == 0:
                        total_text_blocks += 1
                
                # Check for background elements
                if len(page.get_drawings()) > 10:  # Many vector graphics
                    has_complex_background = True
            
            doc.close()
            
            # Decide method based on analysis
            if has_complex_background or total_images > total_text_blocks * 0.3:
                print("[INFO] Detected complex layout/images - using overlay method...")
                return self.translate_pdf_with_overlay(input_path, output_path, target_language)
            else:
                print("[INFO] Detected text-heavy document - using redaction method...")
                return self.translate_pdf_with_redaction(input_path, output_path, target_language)
                
        except Exception as e:
            print(f"Error in hybrid translation: {str(e)}")
            # Fallback to overlay method
            return self.translate_pdf_with_overlay(input_path, output_path, target_language)

    def _int_to_rgb(self, color_int: int) -> tuple:
        """Convert integer color to RGB tuple."""
        if color_int == 0:
            return (0, 0, 0)
        r = (color_int >> 16) & 0xFF
        g = (color_int >> 8) & 0xFF
        b = color_int & 0xFF
        return (r, g, b)

    # ============= Other Methods =============

    def translate_txt(self, input_path: str, output_path: str, 
                     target_language: str) -> bool:
        """
        Translate plain text file.
        
        Args:
            input_path: Path to input text file
            output_path: Path to save translated text file
            target_language: Target language for translation
            
        Returns:
            Success status
        """
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            translated = self.translate_text(text, target_language)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(translated)

            print(f"[SUCCESS] Translated TXT saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating TXT: {str(e)}")
            return False

    def translate_md(self, input_path: str, output_path: str,
                    target_language: str) -> bool:
        """
        Translate Markdown file while preserving formatting, code blocks, and links.

        Args:
            input_path: Path to input Markdown file
            output_path: Path to save translated Markdown file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Split content into blocks (paragraphs, code blocks, etc.)
            lines = content.split('\n')
            blocks = []
            current_block = []
            in_code_block = False
            code_block_delimiter = None

            for line in lines:
                stripped = line.strip()

                # Detect code block boundaries (``` or ~~~)
                is_fence = stripped.startswith('```') or stripped.startswith('~~~')

                if is_fence:
                    if not in_code_block:
                        # Start of code block
                        # Save current text block before code block
                        if current_block:
                            blocks.append(('text', '\n'.join(current_block)))
                            current_block = []
                        in_code_block = True
                        code_block_delimiter = stripped[:3]
                        current_block = [line]
                    else:
                        # End of code block (check if same delimiter)
                        current_block.append(line)
                        blocks.append(('code', '\n'.join(current_block)))
                        current_block = []
                        in_code_block = False
                        code_block_delimiter = None
                elif in_code_block:
                    # Inside code block
                    current_block.append(line)
                else:
                    # Regular text line
                    current_block.append(line)

            # Add remaining block
            if current_block:
                block_type = 'code' if in_code_block else 'text'
                blocks.append((block_type, '\n'.join(current_block)))

            # Prepare blocks for parallel translation
            text_blocks_to_translate = []
            api_key = os.environ.get('GEMINI_API_KEY', None)

            for idx, (block_type, block_content) in enumerate(blocks):
                if block_type == 'text' and block_content.strip():
                    text_blocks_to_translate.append((
                        idx,
                        block_content,
                        target_language,
                        api_key,
                        self.model
                    ))

            # Calculate dynamic worker allocation
            workers_allocated, _ = self._calculate_workers_per_page(max(1, len(text_blocks_to_translate) // 5))

            print(f"Processing Markdown with dynamic allocation:")
            print(f"  - Total blocks: {len(blocks)}")
            print(f"  - Text blocks to translate: {len(text_blocks_to_translate)}")
            print(f"  - Using up to {workers_allocated} workers")

            # Translate text blocks in parallel
            translated_blocks = {}
            if text_blocks_to_translate:
                with ThreadPoolExecutor(max_workers=workers_allocated) as executor:
                    future_to_block = {
                        executor.submit(self.translate_text_block, block_data): block_data[0]
                        for block_data in text_blocks_to_translate
                    }

                    for future in concurrent.futures.as_completed(future_to_block):
                        block_idx = future_to_block[future]
                        try:
                            result = future.result()
                            if 'error' not in result:
                                translated_blocks[block_idx] = result['translated_text']
                        except Exception as e:
                            print(f"Block {block_idx} translation failed: {str(e)}")

            # Reconstruct the markdown with translated text
            translated_content_parts = []
            for idx, (block_type, block_content) in enumerate(blocks):
                if block_type == 'code':
                    # Keep code blocks unchanged
                    translated_content_parts.append(block_content)
                elif idx in translated_blocks:
                    # Use translated text
                    translated_content_parts.append(translated_blocks[idx])
                else:
                    # Keep original (empty blocks or failed translations)
                    translated_content_parts.append(block_content)

            translated_content = '\n'.join(translated_content_parts)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(translated_content)

            print(f"[SUCCESS] Translated Markdown saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating Markdown: {str(e)}")
            return False

    def translate_document(self, input_path: str, output_path: Optional[str] = None,
                          target_language: str = "Spanish", method: str = "auto") -> bool:
        """
        Main method to translate any supported document format.

        Args:
            input_path: Path to input document
            output_path: Path to save translated document (optional)
            target_language: Target language for translation
            method: PDF translation method - "overlay", "redaction", or "auto"

        Returns:
            Success status
        """
        # Determine file type
        file_ext = Path(input_path).suffix.lower()

        # Generate output path if not provided
        if not output_path:
            base = Path(input_path).stem
            dir_path = Path(input_path).parent
            output_path = str(dir_path / f"{base}_translated_{target_language}{file_ext}")

        # Route to appropriate translator
        if file_ext == '.pptx':
            return self.translate_pptx(input_path, output_path, target_language)
        elif file_ext == '.pdf':
            # Choose PDF translation method
            if method == "overlay":
                return self.translate_pdf_with_overlay(input_path, output_path, target_language)
            elif method == "redaction":
                return self.translate_pdf_with_redaction(input_path, output_path, target_language)
            else:  # auto
                return self.translate_pdf_hybrid(input_path, output_path, target_language)
        elif file_ext == '.docx':
            return self.translate_docx(input_path, output_path, target_language)
        elif file_ext in ['.txt', '.text']:
            return self.translate_txt(input_path, output_path, target_language)
        elif file_ext in ['.md', '.markdown']:
            return self.translate_md(input_path, output_path, target_language)
        else:
            print(f"Unsupported file format: {file_ext}")
            return False
    
    def batch_translate(self, input_folder: str, output_folder: str,
                       target_language: str, file_types: List[str] = None) -> Dict:
        """
        Translate multiple documents in a folder.
        
        Args:
            input_folder: Path to folder containing documents
            output_folder: Path to save translated documents
            target_language: Target language for translation
            file_types: List of file extensions to process (default: all supported)
            
        Returns:
            Dictionary with translation results
        """
        if file_types is None:
            file_types = ['.pptx', '.pdf', '.docx', '.txt', '.md', '.markdown']
        
        # Create output folder if it doesn't exist
        Path(output_folder).mkdir(parents=True, exist_ok=True)
        
        results = {'success': [], 'failed': []}

        # Process each file
        for file_path in Path(input_folder).iterdir():
            if file_path.suffix.lower() in file_types:
                print(f"\n[FILE] Processing: {file_path.name}")
                output_path = str(Path(output_folder) / f"{file_path.stem}_translated{file_path.suffix}")

                if self.translate_document(str(file_path), output_path, target_language):
                    results['success'].append(file_path.name)
                else:
                    results['failed'].append(file_path.name)

        # Print summary
        print("\n" + "="*50)
        print("Translation Summary:")
        print(f"[SUCCESS] Successfully translated: {len(results['success'])} files")
        if results['failed']:
            print(f"[FAILED] Failed: {len(results['failed'])} files")
            for file in results['failed']:
                print(f"  - {file}")
        
        return results


# CLI interface
if __name__ == "__main__":


    # Create translator
    translator = DocumentTranslator()

    # Translate document
    translator.translate_document(
        input_path="ex06-questions.pdf",
        target_language="Chinese"
    )
    
    # Example 2: Translate PowerPoint presentation
    # translator.translate_document(
    #     input_path="presentation.pptx",
    #     target_language="French"
    # )
    
    # Example 3: Translate Word document
    # translator.translate_document(
    #     input_path="report.docx",
    #     target_language="Spanish"
    # )

    # Example 4: Translate Markdown file (preserves code blocks, links, formatting)
    # translator.translate_document(
    #     input_path="README.md",
    #     target_language="Chinese"
    # )

    # Example 5: Batch translate all documents
    # translator.batch_translate(
    #     input_folder="documents/",
    #     output_folder="translated/",
    #     target_language="Japanese"
    # )