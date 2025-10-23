"""
PowerPoint (PPTX) translation module.
Preserves formatting, bullet points, and table structures.
"""

import os
import json
from typing import List
from concurrent.futures import ThreadPoolExecutor
import concurrent.futures
from google import genai
from pydantic import BaseModel
from pptx import Presentation

from .core import TranslationClient


class TranslatedRuns(BaseModel):
    """Schema for translated text runs"""
    translations: List[str]


class PPTXTranslator:
    """
    PowerPoint translator that preserves formatting and structure.
    """

    def __init__(self, translation_client: TranslationClient):
        """
        Initialize PPTX translator.

        Args:
            translation_client: Core translation client instance
        """
        self.client = translation_client

    def _translate_runs_with_context(self, runs: List[str], target_language: str,
                                    api_key: str, model: str) -> List[str]:
        """
        Translate runs while preserving context and structure.
        Uses Pydantic schema for guaranteed JSON format.

        Args:
            runs: List of text runs
            target_language: Target language
            api_key: API key
            model: Model name

        Returns:
            List of translated runs
        """
        if not runs:
            return runs

        if api_key:
            os.environ['GEMINI_API_KEY'] = api_key
        genai_client = genai.Client()

        user_prompt = (
            f"You are a professional technical translator. Translate the following text runs into {target_language}.\n\n"
            "Rules:\n"
            f"- Translate each run into {target_language} with native, accurate technical language\n"
            "- Consider all runs together for context, but return the translation per run\n"
            "- Keep the EXACT same number of runs in the same order\n"
            "- Do not merge or split runs; do not drop whitespace\n"
            "- Preserve placeholders, code, variables, URLs, IDs unchanged\n"
            "- Return exactly the same number of translations as input runs\n\n"
            f"Input runs ({len(runs)} items):\n{json.dumps(runs, ensure_ascii=False)}"
        )

        try:
            resp = genai_client.models.generate_content(
                model=model,
                contents=user_prompt,
                config={
                    "response_mime_type": "application/json",
                    "response_schema": TranslatedRuns,
                }
            )

            # Use parsed response for guaranteed structure
            result: TranslatedRuns = resp.parsed

            if len(result.translations) == len(runs):
                return result.translations
            else:
                print(f"Warning: Translation count mismatch. Expected {len(runs)}, got {len(result.translations)}")

        except Exception as e:
            print(f"Error in translation API call: {str(e)}")

        return runs

    def _translate_slide_content(self, slide_data: tuple) -> dict:
        """
        Translate content of a single slide.

        Args:
            slide_data: Tuple containing (slide_index, slide_content, target_language, api_key, model)

        Returns:
            Dictionary with translated content for the slide
        """
        slide_idx, shapes_data, target_language, api_key, model = slide_data

        try:
            if api_key:
                os.environ['GEMINI_API_KEY'] = api_key
            genai_client = genai.Client()
            translated_shapes = []

            for shape_idx, shape_info in enumerate(shapes_data):
                shape_type, content = shape_info

                if shape_type == 'text':
                    if content and content.strip():
                        prompt = (
                            f"You are a professional technical translator. Translate into {target_language} and preserve formatting.\n\n"
                            f"Translate to {target_language} with native, accurate, technical wording.\n"
                            "Strictly preserve original line breaks, indentation, and list markers.\n"
                            "Only return the translated text.\n\n"
                            f"Text to translate:\n{content}"
                        )

                        response = genai_client.models.generate_content(
                            model=model,
                            contents=prompt
                        )
                        translated_shapes.append((shape_idx, response.text.strip()))
                    else:
                        translated_shapes.append((shape_idx, content))

                elif shape_type == 'table':
                    # Translate table content
                    translated_table = []
                    for row_idx, row_cells in enumerate(content):
                        translated_row_cells = []
                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                            translated_cell_paragraphs = []
                            for para_idx, runs in enumerate(cell_paragraphs):
                                if runs and any(r.strip() for r in runs):
                                    translated_runs = self._translate_runs_with_context(
                                        runs, target_language, api_key, model
                                    )
                                else:
                                    translated_runs = runs
                                translated_cell_paragraphs.append(translated_runs)
                            translated_row_cells.append(translated_cell_paragraphs)
                        translated_table.append(translated_row_cells)
                    translated_shapes.append((shape_idx, translated_table))

                elif shape_type == 'paragraphs':
                    translated_paragraphs = []
                    for para_idx, runs_data in enumerate(content):
                        if runs_data and any(r.strip() for r in runs_data):
                            translated_runs = self._translate_runs_with_context(
                                runs_data, target_language, api_key, model
                            )
                        else:
                            translated_runs = runs_data
                        translated_paragraphs.append((
                            para_idx,
                            [(i, t) for i, t in enumerate(translated_runs)]
                        ))
                    translated_shapes.append((shape_idx, translated_paragraphs))

            return {'slide_idx': slide_idx, 'translated_shapes': translated_shapes}

        except Exception as e:
            print(f"Translation error on slide {slide_idx + 1}: {str(e)}")
            return {'slide_idx': slide_idx, 'error': str(e)}

    def translate(self, input_path: str, output_path: str,
                  target_language: str, progress_callback=None) -> bool:
        """
        Translate PowerPoint presentation.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to save translated PPTX file
            target_language: Target language for translation
            progress_callback: Optional callback function(progress: float) for progress updates (0.0 to 1.0)

        Returns:
            Success status
        """
        try:
            if progress_callback:
                progress_callback(0.05)

            prs = Presentation(input_path)

            # Extract content from all slides
            slides_data = []
            for slide_idx, slide in enumerate(prs.slides):
                shapes_data = []

                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_table:
                        table_data = []
                        table = shape.table
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_paragraphs = []
                                try:
                                    paragraphs = list(cell.text_frame.paragraphs)
                                except Exception:
                                    paragraphs = []
                                for paragraph in paragraphs:
                                    runs_data = [run.text for run in paragraph.runs]
                                    cell_paragraphs.append(runs_data)
                                row_cells.append(cell_paragraphs)
                            table_data.append(row_cells)
                        shapes_data.append(('table', table_data))
                    elif shape.has_text_frame:
                        paragraphs_data = []
                        for paragraph in shape.text_frame.paragraphs:
                            runs_data = [run.text for run in paragraph.runs]
                            paragraphs_data.append(runs_data)
                        shapes_data.append(('paragraphs', paragraphs_data))
                    elif hasattr(shape, "text") and shape.text:
                        shapes_data.append(('text', shape.text))
                    else:
                        shapes_data.append(('empty', None))

                if shapes_data:
                    # Use API key from the translation client
                    api_key = self.client.api_key if hasattr(self.client, 'api_key') and self.client.api_key else os.environ.get('GEMINI_API_KEY', None)
                    slides_data.append((slide_idx, shapes_data, target_language, api_key, self.client.model))

            if not slides_data:
                print("No text content found in PPTX")
                prs.save(output_path)
                if progress_callback:
                    progress_callback(1.0)
                return True

            total_slides = len(slides_data)
            workers_per_slide, concurrent_slides = self.client.calculate_workers_per_page(total_slides)

            print(f"Processing {total_slides} slides with dynamic allocation:")
            print(f"  - Processing {concurrent_slides} slides concurrently")
            print(f"  - Using up to {workers_per_slide} workers per slide")
            print(f"  - Total workers: {workers_per_slide * concurrent_slides}")

            if progress_callback:
                progress_callback(0.1)

            # Process slides in parallel
            translated_results = {}
            completed_slides = 0
            with ThreadPoolExecutor(max_workers=workers_per_slide * concurrent_slides) as executor:
                future_to_slide = {
                    executor.submit(self._translate_slide_content, slide_data): slide_data[0]
                    for slide_data in slides_data
                }

                for future in concurrent.futures.as_completed(future_to_slide):
                    slide_idx = future_to_slide[future]
                    try:
                        result = future.result()
                        if 'error' not in result:
                            translated_results[slide_idx] = result['translated_shapes']
                        completed_slides += 1
                        # Report progress from 10% to 80% during translation
                        if progress_callback:
                            progress = 0.1 + (completed_slides / total_slides) * 0.7
                            progress_callback(progress)
                    except Exception as e:
                        print(f"Slide {slide_idx + 1} translation failed: {str(e)}")
                        completed_slides += 1
                        if progress_callback:
                            progress = 0.1 + (completed_slides / total_slides) * 0.7
                            progress_callback(progress)

            # Apply translations back to the presentation
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx in translated_results:
                    shape_translations = translated_results[slide_idx]
                    shape_translations.sort(key=lambda x: x[0])

                    for shape_trans_idx, (shape_idx, translated_content) in enumerate(shape_translations):
                        if shape_idx < len(list(slide.shapes)):
                            shape = list(slide.shapes)[shape_idx]

                            if shape.has_text_frame and isinstance(translated_content, list):
                                for para_trans in translated_content:
                                    if isinstance(para_trans, tuple) and len(para_trans) == 2:
                                        para_idx, translated_runs = para_trans
                                        if para_idx < len(list(shape.text_frame.paragraphs)):
                                            paragraph = list(shape.text_frame.paragraphs)[para_idx]
                                            for run_trans in translated_runs:
                                                if isinstance(run_trans, tuple) and len(run_trans) == 2:
                                                    run_idx, run_text = run_trans
                                                    if run_idx < len(list(paragraph.runs)):
                                                        list(paragraph.runs)[run_idx].text = run_text
                            elif shape.has_table and isinstance(translated_content, list):
                                table = shape.table
                                for row_idx, row_cells in enumerate(translated_content):
                                    if row_idx < len(list(table.rows)):
                                        row = list(table.rows)[row_idx]
                                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                                            if cell_idx < len(list(row.cells)):
                                                cell = list(row.cells)[cell_idx]
                                                for para_idx, translated_runs in enumerate(cell_paragraphs):
                                                    if para_idx < len(list(cell.text_frame.paragraphs)):
                                                        paragraph = list(cell.text_frame.paragraphs)[para_idx]
                                                        for run_idx, run_text in enumerate(translated_runs):
                                                            if run_idx < len(list(paragraph.runs)):
                                                                list(paragraph.runs)[run_idx].text = run_text
                            elif hasattr(shape, "text") and isinstance(translated_content, str):
                                shape.text = translated_content

            # Save translated presentation
            if progress_callback:
                progress_callback(0.9)

            prs.save(output_path)

            if progress_callback:
                progress_callback(1.0)

            print(f"[SUCCESS] Translated PPTX saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating PPTX: {str(e)}")
            return False
