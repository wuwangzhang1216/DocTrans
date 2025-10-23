"""
PowerPoint (PPTX) translation module.
Preserves formatting, bullet points, and table structures.

Architecture:
- Uses nested parallelization for optimal performance
- Outer level: Multiple slides processed concurrently
- Inner level: All text elements within each slide processed in parallel
- Flattens hierarchical structure (shapes -> paragraphs -> runs) into parallel tasks
- Reconstructs structure after translation using metadata tracking
"""

import os
import json
from typing import List
from concurrent.futures import ThreadPoolExecutor
import concurrent.futures
from google import genai
from pydantic import BaseModel
from pptx import Presentation

from .core import TranslationClient


class TranslatedRuns(BaseModel):
    """Schema for translated text runs"""
    translations: List[str]


class PPTXTranslator:
    """
    PowerPoint translator that preserves formatting and structure.
    """

    def __init__(self, translation_client: TranslationClient):
        """
        Initialize PPTX translator.

        Args:
            translation_client: Core translation client instance
        """
        self.client = translation_client

    def _translate_runs_with_context(self, runs: List[str], target_language: str,
                                    api_key: str, model: str) -> List[str]:
        """
        Translate runs while preserving context and structure.
        Uses Pydantic schema for guaranteed JSON format.

        Args:
            runs: List of text runs
            target_language: Target language
            api_key: API key
            model: Model name

        Returns:
            List of translated runs
        """
        if not runs:
            return runs

        if api_key:
            os.environ['GEMINI_API_KEY'] = api_key
        genai_client = genai.Client()

        user_prompt = (
            f"You are a professional technical translator. Translate the following text runs into {target_language}.\n\n"
            "Rules:\n"
            f"- Translate each run into {target_language} with native, accurate technical language\n"
            "- Consider all runs together for context, but return the translation per run\n"
            "- Keep the EXACT same number of runs in the same order\n"
            "- Do not merge or split runs; do not drop whitespace\n"
            "- Preserve placeholders, code, variables, URLs, IDs unchanged\n"
            "- Return exactly the same number of translations as input runs\n\n"
            f"Input runs ({len(runs)} items):\n{json.dumps(runs, ensure_ascii=False)}"
        )

        try:
            resp = genai_client.models.generate_content(
                model=model,
                contents=user_prompt,
                config={
                    "response_mime_type": "application/json",
                    "response_schema": TranslatedRuns,
                }
            )

            # Use parsed response for guaranteed structure
            result: TranslatedRuns = resp.parsed

            if len(result.translations) == len(runs):
                return result.translations
            else:
                print(f"Warning: Translation count mismatch. Expected {len(runs)}, got {len(result.translations)}")

        except Exception as e:
            print(f"Error in translation API call: {str(e)}")

        return runs

    def _translate_single_text(self, text: str, target_language: str,
                              api_key: str, model: str) -> str:
        """
        Translate a single text block.

        Args:
            text: Text to translate
            target_language: Target language
            api_key: API key
            model: Model name

        Returns:
            Translated text
        """
        if api_key:
            os.environ['GEMINI_API_KEY'] = api_key
        genai_client = genai.Client()

        prompt = (
            f"You are a professional technical translator. Translate into {target_language} and preserve formatting.\n\n"
            f"Translate to {target_language} with native, accurate, technical wording.\n"
            "Strictly preserve original line breaks, indentation, and list markers.\n"
            "Only return the translated text.\n\n"
            f"Text to translate:\n{text}"
        )

        response = genai_client.models.generate_content(
            model=model,
            contents=prompt
        )
        return response.text.strip()

    def _process_translation_task(self, task: dict) -> dict:
        """
        Process a single translation task (either text or runs).

        Args:
            task: Dictionary containing task type, content, metadata, and translation params

        Returns:
            Dictionary with translated content and metadata for reconstruction
        """
        try:
            task_type = task['task_type']
            content = task['content']
            target_language = task['target_language']
            api_key = task['api_key']
            model = task['model']
            metadata = task['metadata']

            if task_type == 'text':
                translated_content = self._translate_single_text(
                    content, target_language, api_key, model
                )
            elif task_type == 'runs':
                translated_content = self._translate_runs_with_context(
                    content, target_language, api_key, model
                )
            else:
                translated_content = content

            return {
                'success': True,
                'content': translated_content,
                'metadata': metadata
            }

        except Exception as e:
            print(f"Translation task error: {str(e)}")
            return {
                'success': False,
                'content': content,
                'metadata': metadata,
                'error': str(e)
            }

    def _reconstruct_slide_structure(self, translation_results: List[dict],
                                    shapes_data: List[tuple]) -> List[tuple]:
        """
        Reconstruct hierarchical slide structure from flat translation results.

        Args:
            translation_results: Flat list of translation results with metadata
            shapes_data: Original shapes data for structure reference

        Returns:
            List of (shape_idx, translated_content) tuples in original format
        """
        # Group results by shape_idx and shape_type
        shape_groups = {}
        for result in translation_results:
            metadata = result['metadata']
            shape_idx = metadata['shape_idx']
            shape_type = metadata['shape_type']

            if shape_idx not in shape_groups:
                shape_groups[shape_idx] = {
                    'shape_type': shape_type,
                    'results': []
                }
            shape_groups[shape_idx]['results'].append(result)

        # Reconstruct each shape
        translated_shapes = []
        for shape_idx in sorted(shape_groups.keys()):
            group = shape_groups[shape_idx]
            shape_type = group['shape_type']
            results = group['results']

            if shape_type == 'text':
                # Simple text shape - should only have one result
                if results:
                    translated_shapes.append((shape_idx, results[0]['content']))

            elif shape_type == 'paragraphs':
                # Reconstruct paragraphs structure
                para_dict = {}
                for result in results:
                    para_idx = result['metadata']['para_idx']
                    translated_runs = result['content']
                    para_dict[para_idx] = translated_runs

                # Build paragraphs list with indexed runs
                translated_paragraphs = []
                for para_idx in sorted(para_dict.keys()):
                    translated_runs = para_dict[para_idx]
                    translated_paragraphs.append((
                        para_idx,
                        [(i, t) for i, t in enumerate(translated_runs)]
                    ))
                translated_shapes.append((shape_idx, translated_paragraphs))

            elif shape_type == 'table':
                # Reconstruct table structure: row -> cell -> paragraph
                # Get original table structure for dimensions
                original_shape = shapes_data[shape_idx]
                original_table = original_shape[1]

                # Build nested structure
                table_dict = {}
                for result in results:
                    row_idx = result['metadata']['row_idx']
                    cell_idx = result['metadata']['cell_idx']
                    para_idx = result['metadata']['para_idx']
                    translated_runs = result['content']

                    if row_idx not in table_dict:
                        table_dict[row_idx] = {}
                    if cell_idx not in table_dict[row_idx]:
                        table_dict[row_idx][cell_idx] = {}
                    table_dict[row_idx][cell_idx][para_idx] = translated_runs

                # Rebuild table matching original structure
                translated_table = []
                for row_idx in range(len(original_table)):
                    translated_row = []
                    row_data = original_table[row_idx]
                    for cell_idx in range(len(row_data)):
                        translated_cell_paras = []
                        cell_data = original_table[row_idx][cell_idx]
                        for para_idx in range(len(cell_data)):
                            if (row_idx in table_dict and
                                cell_idx in table_dict[row_idx] and
                                para_idx in table_dict[row_idx][cell_idx]):
                                translated_runs = table_dict[row_idx][cell_idx][para_idx]
                            else:
                                # Use original if not translated (empty content)
                                translated_runs = original_table[row_idx][cell_idx][para_idx]
                            translated_cell_paras.append(translated_runs)
                        translated_row.append(translated_cell_paras)
                    translated_table.append(translated_row)

                translated_shapes.append((shape_idx, translated_table))

        return translated_shapes

    def _translate_slide_content(self, slide_data: tuple) -> dict:
        """
        Translate content of a single slide using parallel processing.

        Args:
            slide_data: Tuple containing (slide_index, slide_content, target_language, api_key, model, workers_per_slide)

        Returns:
            Dictionary with translated content for the slide
        """
        slide_idx, shapes_data, target_language, api_key, model, workers_per_slide = slide_data

        try:
            # Build flat list of translation tasks with metadata
            translation_tasks = []

            for shape_idx, shape_info in enumerate(shapes_data):
                shape_type, content = shape_info

                if shape_type == 'text':
                    if content and content.strip():
                        translation_tasks.append({
                            'task_type': 'text',
                            'content': content,
                            'target_language': target_language,
                            'api_key': api_key,
                            'model': model,
                            'metadata': {
                                'shape_idx': shape_idx,
                                'shape_type': 'text'
                            }
                        })
                    else:
                        # Empty text, no translation needed
                        translation_tasks.append({
                            'task_type': 'skip',
                            'content': content,
                            'target_language': target_language,
                            'api_key': api_key,
                            'model': model,
                            'metadata': {
                                'shape_idx': shape_idx,
                                'shape_type': 'text'
                            }
                        })

                elif shape_type == 'table':
                    # Flatten table structure into individual paragraph translation tasks
                    for row_idx, row_cells in enumerate(content):
                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                            for para_idx, runs in enumerate(cell_paragraphs):
                                if runs and any(r.strip() for r in runs):
                                    translation_tasks.append({
                                        'task_type': 'runs',
                                        'content': runs,
                                        'target_language': target_language,
                                        'api_key': api_key,
                                        'model': model,
                                        'metadata': {
                                            'shape_idx': shape_idx,
                                            'shape_type': 'table',
                                            'row_idx': row_idx,
                                            'cell_idx': cell_idx,
                                            'para_idx': para_idx
                                        }
                                    })
                                else:
                                    translation_tasks.append({
                                        'task_type': 'skip',
                                        'content': runs,
                                        'target_language': target_language,
                                        'api_key': api_key,
                                        'model': model,
                                        'metadata': {
                                            'shape_idx': shape_idx,
                                            'shape_type': 'table',
                                            'row_idx': row_idx,
                                            'cell_idx': cell_idx,
                                            'para_idx': para_idx
                                        }
                                    })

                elif shape_type == 'paragraphs':
                    # Flatten paragraphs into individual translation tasks
                    for para_idx, runs_data in enumerate(content):
                        if runs_data and any(r.strip() for r in runs_data):
                            translation_tasks.append({
                                'task_type': 'runs',
                                'content': runs_data,
                                'target_language': target_language,
                                'api_key': api_key,
                                'model': model,
                                'metadata': {
                                    'shape_idx': shape_idx,
                                    'shape_type': 'paragraphs',
                                    'para_idx': para_idx
                                }
                            })
                        else:
                            translation_tasks.append({
                                'task_type': 'skip',
                                'content': runs_data,
                                'target_language': target_language,
                                'api_key': api_key,
                                'model': model,
                                'metadata': {
                                    'shape_idx': shape_idx,
                                    'shape_type': 'paragraphs',
                                    'para_idx': para_idx
                                }
                            })

            # Process all translation tasks in parallel
            translation_results = []
            if translation_tasks:
                with ThreadPoolExecutor(max_workers=workers_per_slide) as executor:
                    future_to_task = {
                        executor.submit(self._process_translation_task, task): task
                        for task in translation_tasks
                    }

                    for future in concurrent.futures.as_completed(future_to_task):
                        try:
                            result = future.result()
                            translation_results.append(result)
                        except Exception as e:
                            task = future_to_task[future]
                            print(f"Task execution error: {str(e)}")
                            translation_results.append({
                                'success': False,
                                'content': task['content'],
                                'metadata': task['metadata'],
                                'error': str(e)
                            })

            # Reconstruct hierarchical structure from flat results
            translated_shapes = self._reconstruct_slide_structure(
                translation_results, shapes_data
            )

            return {'slide_idx': slide_idx, 'translated_shapes': translated_shapes}

        except Exception as e:
            print(f"Translation error on slide {slide_idx + 1}: {str(e)}")
            return {'slide_idx': slide_idx, 'error': str(e)}

    def translate(self, input_path: str, output_path: str,
                  target_language: str) -> bool:
        """
        Translate PowerPoint presentation.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to save translated PPTX file
            target_language: Target language for translation

        Returns:
            Success status
        """
        try:
            prs = Presentation(input_path)

            # Extract content from all slides
            slides_data = []
            for slide_idx, slide in enumerate(prs.slides):
                shapes_data = []

                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_table:
                        table_data = []
                        table = shape.table
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_paragraphs = []
                                try:
                                    paragraphs = list(cell.text_frame.paragraphs)
                                except Exception:
                                    paragraphs = []
                                for paragraph in paragraphs:
                                    runs_data = [run.text for run in paragraph.runs]
                                    cell_paragraphs.append(runs_data)
                                row_cells.append(cell_paragraphs)
                            table_data.append(row_cells)
                        shapes_data.append(('table', table_data))
                    elif shape.has_text_frame:
                        paragraphs_data = []
                        for paragraph in shape.text_frame.paragraphs:
                            runs_data = [run.text for run in paragraph.runs]
                            paragraphs_data.append(runs_data)
                        shapes_data.append(('paragraphs', paragraphs_data))
                    elif hasattr(shape, "text") and shape.text:
                        shapes_data.append(('text', shape.text))
                    else:
                        shapes_data.append(('empty', None))

                if shapes_data:
                    # Use API key from the translation client
                    api_key = self.client.api_key if hasattr(self.client, 'api_key') and self.client.api_key else os.environ.get('GEMINI_API_KEY', None)
                    # Note: workers_per_slide will be added after calculation
                    slides_data.append((slide_idx, shapes_data, target_language, api_key, self.client.model))

            if not slides_data:
                print("No text content found in PPTX")
                prs.save(output_path)
                return True

            total_slides = len(slides_data)
            workers_per_slide, concurrent_slides = self.client.calculate_workers_per_page(total_slides)

            print(f"Processing {total_slides} slides with nested parallel processing:")
            print(f"  - Processing {concurrent_slides} slides concurrently")
            print(f"  - Using up to {workers_per_slide} workers per slide for parallel element translation")
            print(f"  - Total workers: {workers_per_slide * concurrent_slides}")

            # Add workers_per_slide to each slide data tuple
            slides_data_with_workers = [
                (slide_idx, shapes_data, target_language, api_key, model, workers_per_slide)
                for slide_idx, shapes_data, target_language, api_key, model in slides_data
            ]

            # Process slides in parallel
            translated_results = {}
            with ThreadPoolExecutor(max_workers=concurrent_slides) as executor:
                future_to_slide = {
                    executor.submit(self._translate_slide_content, slide_data): slide_data[0]
                    for slide_data in slides_data_with_workers
                }

                for future in concurrent.futures.as_completed(future_to_slide):
                    slide_idx = future_to_slide[future]
                    try:
                        result = future.result()
                        if 'error' not in result:
                            translated_results[slide_idx] = result['translated_shapes']
                    except Exception as e:
                        print(f"Slide {slide_idx + 1} translation failed: {str(e)}")

            # Apply translations back to the presentation
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx in translated_results:
                    shape_translations = translated_results[slide_idx]
                    shape_translations.sort(key=lambda x: x[0])

                    for shape_trans_idx, (shape_idx, translated_content) in enumerate(shape_translations):
                        if shape_idx < len(list(slide.shapes)):
                            shape = list(slide.shapes)[shape_idx]

                            if shape.has_text_frame and isinstance(translated_content, list):
                                for para_trans in translated_content:
                                    if isinstance(para_trans, tuple) and len(para_trans) == 2:
                                        para_idx, translated_runs = para_trans
                                        if para_idx < len(list(shape.text_frame.paragraphs)):
                                            paragraph = list(shape.text_frame.paragraphs)[para_idx]
                                            for run_trans in translated_runs:
                                                if isinstance(run_trans, tuple) and len(run_trans) == 2:
                                                    run_idx, run_text = run_trans
                                                    if run_idx < len(list(paragraph.runs)):
                                                        list(paragraph.runs)[run_idx].text = run_text
                            elif shape.has_table and isinstance(translated_content, list):
                                table = shape.table
                                for row_idx, row_cells in enumerate(translated_content):
                                    if row_idx < len(list(table.rows)):
                                        row = list(table.rows)[row_idx]
                                        for cell_idx, cell_paragraphs in enumerate(row_cells):
                                            if cell_idx < len(list(row.cells)):
                                                cell = list(row.cells)[cell_idx]
                                                for para_idx, translated_runs in enumerate(cell_paragraphs):
                                                    if para_idx < len(list(cell.text_frame.paragraphs)):
                                                        paragraph = list(cell.text_frame.paragraphs)[para_idx]
                                                        for run_idx, run_text in enumerate(translated_runs):
                                                            if run_idx < len(list(paragraph.runs)):
                                                                list(paragraph.runs)[run_idx].text = run_text
                            elif hasattr(shape, "text") and isinstance(translated_content, str):
                                shape.text = translated_content

            # Save translated presentation
            prs.save(output_path)
            print(f"[SUCCESS] Translated PPTX saved to: {output_path}")
            return True

        except Exception as e:
            print(f"Error translating PPTX: {str(e)}")
            return False
